import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
from io import BytesIO

# Colors and style
PRIMARY_BLUE = RGBColor(17, 85, 204)   # Google blue
LIGHT_GREY = RGBColor(240, 240, 240)   # Very light grey for shape backgrounds
DARK_GREY = RGBColor(100, 100, 100)
TITLE_FONT = 'Arial'
BODY_FONT = 'Calibri'

# Helper to add a blue section title bar
def add_section_bar(slide, title, top=0.15):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(13.33), Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.width = Pt(0)
    tf = bar.text_frame
    tf.text = "  " + title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    tf.paragraphs[0].font.name = TITLE_FONT
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    bar.shadow.inherit = False

# Helper to add content block
def add_content_block(slide, text, top=1.1, height=2.3, font_pt=28, bullet=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(12.3), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.width = Pt(0)
    tf = shape.text_frame
    lines = text.split('\n')
    # Set the default paragraph only if appropriate, else add paragraphs
    if len(lines) == 1:
        tf.paragraphs[0].text = lines[0]
        tf.paragraphs[0].font.size = Pt(font_pt)
        tf.paragraphs[0].font.name = BODY_FONT
        if bullet:
            tf.paragraphs[0].level = 0
    else:
        tf.clear()  # Remove default empty, so there's nothing to delete
        for para in lines:
            p = tf.add_paragraph()
            p.text = para
            p.font.size = Pt(font_pt)
            p.font.name = BODY_FONT
            if bullet:
                p.level = 0
    return shape

# Helper formula slide
def add_formula_slide(prs, section, title, formula, description):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bar(slide, section)
    # Big formula
    fig, ax = plt.subplots(figsize=(8, 1.5))
    ax.text(0.5, 0.5, f'${formula}$', fontsize=48, ha='center', va='center')
    ax.axis('off')
    buf = BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True, dpi=200)
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1.5), Inches(1.25), width=Inches(10))
    # Description
    add_content_block(slide, description, top=3.2, height=1.6, font_pt=28, bullet=False)
    # Note for image placement
    slide.notes_slide.notes_text_frame.text = 'Formula image generated. Add figure/chart below if desired.'

# Slide layout
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

######################
# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Medical Insurance Cost Prediction"
slide.placeholders[1].text = "Machine Learning Project\nYour Name\nNov 2025"

######################
# Section: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Project Overview")
add_content_block(slide, (
    "- Predict insurance charges using regression models\n"
    "- Data analysis and ML pipeline (EDA, preprocessing, modeling, evaluation)\n"
    "- Compare Linear, Ridge, and Lasso regression"),
    top=1.1, height=2.3)
slide.notes_slide.notes_text_frame.text = 'Use relevant iconography or illustration for topic if available.'

######################
# Section: Dataset
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Dataset and Features")
add_content_block(slide, (
    "- Synthetic dataset mimicking medical insurance data\n"
    "  (1,338 records, 7 columns: age, sex, bmi, children, smoker, region, charges)\n"
    "- Target: charges (medical cost)",
    ), top=1.1, height=1.8)
# Add a numbered block sample
sample_txt = (
    " 1. \t age: integer, patient age\n"
    " 2. \t sex: male/female\n"
    " 3. \t bmi: float, body mass index\n"
    " 4. \t children: int, dependents\n"
    " 5. \t smoker: yes/no\n"
    " 6. \t region: string\n"
    " 7. \t charges: float, target variable\n"
)
add_content_block(slide, sample_txt, top=3, height=2.0, font_pt=20, bullet=False)
slide.notes_slide.notes_text_frame.text = 'Show head() or describe() output as a screenshot below.'

######################
# Section: Exploratory Data Analysis
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Exploratory Data Analysis (EDA)")
add_content_block(slide, (
    "- Correlation heatmap for numerical features\n"
    "- Pairplot for feature relationships\n"
    "- High correlation noted with smoker, age, bmi and charges"
    ), top=1.1, height=1.8)
slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(12), Inches(1.8)).text = (
        "[Add: Correlation heatmap and pairplot image here]")
slide.notes_slide.notes_text_frame.text = 'Paste notebook’s EDA charts in this region.'

######################
# Section: Preprocessing Steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Preprocessing")
add_content_block(slide, (
    "- One-hot encoding for categorical features\n"
    "- Standard scaling (mean=0, var=1)\n"
    "- Ensures ML models handle features optimally"
    ), top=1.1, height=2.0)
slide.notes_slide.notes_text_frame.text = 'Expand if adding pipeline block diagram.'

######################
# Section: Modeling
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Modeling: Train/Test Split & Models")
add_content_block(slide, (
    "Step 1: 80/20 train-test split\n"
    "Step 2: Train Linear Regression\n"
    "Step 3: Train Ridge, Lasso (α tuned for regularization)\n"
    "Step 4: Evaluate and compare"
    ), top=1.1, height=2.3, font_pt=26, bullet=False)

# Model formulas
add_formula_slide(
    prs,
    "Modeling: Formula",
    "Linear Regression",
    r"\hat{y} = w_0 + w_1x_1 + ... + w_nx_n",
    "Predicts charges as a weighted sum. Simple, interpretable baseline model."
)
add_formula_slide(
    prs,
    "Modeling: Formula",
    "Ridge Regression (L2)",
    r"\mathrm{Cost} = \sum_{i=1}^{n}(y_i - \hat{y}_i)^2 + \alpha \sum_{j=1}^{p}w_j^2",
    "Adds penalty for large weights (L2). Controls overfitting, improves generalization."
)
add_formula_slide(
    prs,
    "Modeling: Formula",
    "Lasso Regression (L1)",
    r"\mathrm{Cost} = \sum_{i=1}^{n}(y_i - \hat{y}_i)^2 + \alpha \sum_{j=1}^{p}|w_j|",
    "Penalty for sum of absolute weights. Can shrink coefficients to zero (feature selection)."
)

######################
# Section: Evaluation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Model Evaluation Metrics")
add_content_block(slide, (
    "- RMSE (Root Mean Squared Error): Lower = Better\n"
    "- R² (R squared score): 1 = perfect, 0 = no fit\n"
    "- Linear: RMSE ~ 3,153, R² ~ 0.88\n"
    "- Ridge: RMSE ~ 3,152, R² ~ 0.88\n"
    "- Lasso: RMSE ~ 3,153, R² ~ 0.88"
    ), top=1.1, height=2.15)
slide.notes_slide.notes_text_frame.text = 'Insert metrics table or screenshot/plot here.'

add_formula_slide(
    prs, "Evaluation", "Key Metric Formulas",
    r"\mathrm{RMSE} = \sqrt{\frac{1}{n}\sum_{i=1}^{n}(y_i-\hat{y}_i)^2}\qquad R^2 = 1-\frac{SS_{res}}{SS_{tot}}",
    "Lower RMSE = better fit. R² close to 1 = good model performance."
)

######################
# Section: Feature Importance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Feature Importance")
add_content_block(slide, (
    "- Largest coefficients = most influence on prediction\n"
    "- Smoking, age, BMI had highest weights\n"
    ), top=1.1, height=1.3)
slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(1)).text = (
    "[Add: Coefficient table screenshot here]")
slide.notes_slide.notes_text_frame.text = 'Paste coefficients table for interpretation.'

######################
# Section: Residuals
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Residual Analysis")
add_content_block(slide, (
    "- Residual = Actual - Predicted\n"
    "- Most residuals close to zero: Good model fit\n"
    ), top=1.1, height=1.2)
slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(1)).text = (
    "[Add: Residual plot screenshot here]")
slide.notes_slide.notes_text_frame.text = 'Paste residual plot for visual fit.'

######################
# Section: Prediction Example
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Sample Prediction")
add_content_block(slide, (
    "Input example: age=40, bmi=28, children=2, male, smoker, southeast\n"
    "Model predicts charges: ₹ 41,598.94\n"
    "\n[Insert output from custom test case]"
    ), top=1.1, height=1.8)
slide.notes_slide.notes_text_frame.text = 'Copy notebook output for custom test.'

######################
# Section: Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Conclusions")
add_content_block(slide, (
    "- Regression models offer reliable cost prediction\n"
    "- Largest impact: smoking, age, BMI\n"
    "- Extendable pipeline, baseline for production"
), top=1.1, height=1.2)

######################
# Section: Future Work
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bar(slide, "Future Work")
add_content_block(slide, (
    "- More advanced regressors (Random Forest, XGBoost)\n"
    "- Hyperparameter tuning\n"
    "- Richer features, real-world data"
    ), top=1.1, height=1.3)

######################
# Section: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Questions?"
slide.placeholders[1].text = "Thank you! Open for discussion."

prs.save("insurance_cost_prediction_presentation.pptx")
print("Modern Google Slides-style PPTX Saved: insurance_cost_prediction_presentation.pptx")
